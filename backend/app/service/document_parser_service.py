# Copyright © 2026 深圳市深维智见教育科技有限公司 版权所有
# 未经授权，禁止转售或仿制。

"""
统一文档解析服务（本地优先 + DocMind 兜底）

背景（修复问题1/问题2的核心）：
1. 旧实现把所有文件（包括 .md/.txt 这类纯文本）都推给阿里云 DocMind，
   而当前账号未开通 DocMind（接口返回 HTTP 200 + Code=DocMindServiceNotOpen），
   SDK 不抛异常，导致 submit_job() 返回 None，最终只留下 "文档提交失败" 这种无信息量的错误。
2. 聊天附件的 PDF/Word 只写入 "[PDF 文件: xxx]" 占位符，RAG 无法引用真实内容。

本模块提供唯一入口 parse_document()：
- 纯文本/代码/表格类：本地直接解析，不依赖任何外部 API
- PDF/Word/Excel/PPT：优先本地库解析，失败或缺库时才回退 DocMind
- 图片：只能走 DocMind OCR
- 任何失败都抛 DocumentParseError，且 message 里带真实原因（含 DocMind 的 Code/Message）
"""

import json
import logging
import os
from dataclasses import dataclass, field
from typing import List, Optional

logger = logging.getLogger(__name__)


class DocumentParseError(Exception):
    """文档解析失败（携带真实原因）"""


# ==================== 扩展名分类 ====================

# 纯文本类：直接按文本读取
PLAIN_TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".log", ".rst",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".go", ".c", ".h", ".cpp",
    ".sh", ".sql", ".ini", ".conf", ".toml", ".yaml", ".yml",
}

# 结构化文本类：需要轻量处理
JSON_EXTENSIONS = {".json"}
CSV_EXTENSIONS = {".csv", ".tsv"}
HTML_EXTENSIONS = {".html", ".htm", ".xml"}

# 二进制文档类：本地库优先，DocMind 兜底
PDF_EXTENSIONS = {".pdf"}
WORD_EXTENSIONS = {".docx", ".doc"}
EXCEL_EXTENSIONS = {".xlsx", ".xls"}
PPT_EXTENSIONS = {".pptx", ".ppt"}

# 图片类：只能 OCR（DocMind）
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".tiff"}


@dataclass
class ParsedDocument:
    """解析结果"""
    text: str
    parser: str                       # 实际使用的解析器名称，便于排障
    warnings: List[str] = field(default_factory=list)

    @property
    def char_count(self) -> int:
        return len(self.text)


def get_file_extension(filename: str) -> str:
    return os.path.splitext(filename or "")[1].lower()


def is_text_like(filename: str) -> bool:
    """是否属于可以本地纯文本解析的类型"""
    ext = get_file_extension(filename)
    return (
        ext in PLAIN_TEXT_EXTENSIONS
        or ext in JSON_EXTENSIONS
        or ext in CSV_EXTENSIONS
        or ext in HTML_EXTENSIONS
    )


# ==================== 本地解析实现 ====================

def _read_text_file(file_path: str) -> str:
    """按多种编码尝试读取文本文件"""
    last_error: Optional[Exception] = None
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "latin-1"):
        try:
            with open(file_path, "r", encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError as e:
            last_error = e
            continue
    # 最后兜底：忽略错误
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:  # pragma: no cover
        raise DocumentParseError(f"文本文件读取失败: {e}") from (last_error or e)


def _parse_json(file_path: str) -> str:
    raw = _read_text_file(file_path)
    try:
        data = json.loads(raw)
        return json.dumps(data, ensure_ascii=False, indent=2)
    except Exception:
        # JSON 非法时按纯文本处理，不阻断入库
        return raw


def _parse_html(file_path: str) -> str:
    raw = _read_text_file(file_path)
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return raw
    soup = BeautifulSoup(raw, "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n")


def _parse_csv(file_path: str) -> str:
    """CSV/TSV 转成可读文本（保留表头语义）"""
    import csv as _csv

    raw = _read_text_file(file_path)
    delimiter = "\t" if get_file_extension(file_path) == ".tsv" else ","
    lines: List[str] = []
    try:
        reader = _csv.reader(raw.splitlines(), delimiter=delimiter)
        rows = list(reader)
    except Exception:
        return raw

    if not rows:
        return ""

    header = rows[0]
    lines.append(" | ".join(header))
    for row in rows[1:]:
        if not any(cell.strip() for cell in row):
            continue
        if len(row) == len(header):
            lines.append(
                "; ".join(
                    f"{h.strip()}: {c.strip()}" for h, c in zip(header, row) if c.strip()
                )
            )
        else:
            lines.append(" | ".join(row))
    return "\n".join(lines)


def _parse_pdf_local(file_path: str) -> str:
    """本地 PDF 解析（pypdf / PyPDF2 任一可用即可）"""
    reader = None
    try:
        from pypdf import PdfReader  # type: ignore

        reader = PdfReader(file_path)
    except ImportError:
        try:
            from PyPDF2 import PdfReader  # type: ignore

            reader = PdfReader(file_path)
        except ImportError as e:
            raise DocumentParseError(
                "本地 PDF 解析库不可用（请安装 pypdf 或 PyPDF2）"
            ) from e

    try:
        pages: List[str] = []
        for index, page in enumerate(reader.pages):
            try:
                pages.append(page.extract_text() or "")
            except Exception as e:  # 单页失败不影响整体
                logger.warning("PDF 第 %s 页解析失败: %s", index + 1, e)
        return "\n\n".join(p for p in pages if p.strip())
    except Exception as e:
        raise DocumentParseError(f"本地 PDF 解析失败: {type(e).__name__}: {e}") from e


def _parse_docx_local(file_path: str) -> str:
    try:
        import docx  # type: ignore
    except ImportError as e:
        raise DocumentParseError("本地 Word 解析库不可用（请安装 python-docx）") from e

    try:
        document = docx.Document(file_path)
    except Exception as e:
        raise DocumentParseError(f"本地 Word 解析失败: {type(e).__name__}: {e}") from e

    parts = [p.text for p in document.paragraphs if p.text and p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _parse_xlsx_local(file_path: str) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError as e:
        raise DocumentParseError("本地 Excel 解析库不可用（请安装 openpyxl）") from e

    try:
        wb = load_workbook(file_path, read_only=True, data_only=True)
    except Exception as e:
        raise DocumentParseError(f"本地 Excel 解析失败: {type(e).__name__}: {e}") from e

    parts: List[str] = []
    for ws in wb.worksheets:
        parts.append(f"# 工作表: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)


def _parse_pptx_local(file_path: str) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as e:
        raise DocumentParseError("本地 PPT 解析库不可用（请安装 python-pptx）") from e

    try:
        prs = Presentation(file_path)
    except Exception as e:
        raise DocumentParseError(f"本地 PPT 解析失败: {type(e).__name__}: {e}") from e

    parts: List[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"# 第 {i} 页")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text.strip())
    return "\n".join(parts)


# ==================== DocMind 兜底 ====================

def _parse_with_docmind(file_path: str, file_name: str) -> str:
    """
    使用 DocMind 解析（仅在本地解析不可用时调用）

    未配置或服务未开通时抛 DocumentParseError，并带上阿里云返回的真实 Code/Message。
    """
    from service.docmind_service import DocMindError, DocMindService, is_docmind_configured

    if not is_docmind_configured():
        raise DocumentParseError(
            "DocMind 未配置（缺少 DOCMIND_ACCESS_KEY_ID/DOCMIND_ACCESS_KEY_SECRET），"
            "且该文件类型无法本地解析"
        )

    try:
        service = DocMindService()
        task_id = service.submit_job(file_path, file_name)
        service.wait_for_completion(task_id)
        return service.collect_all_results(task_id)
    except DocMindError as e:
        raise DocumentParseError(f"DocMind 解析失败: {e}") from e
    except Exception as e:
        raise DocumentParseError(
            f"DocMind 解析异常: {type(e).__name__}: {e}"
        ) from e


# ==================== 统一入口 ====================

def parse_document(
    file_path: str,
    file_name: Optional[str] = None,
    allow_docmind: bool = True,
) -> ParsedDocument:
    """
    解析文档为纯文本

    Args:
        file_path: 本地文件路径
        file_name: 原始文件名（用于判断类型，默认取 file_path 的 basename）
        allow_docmind: 本地解析不可用时是否允许回退 DocMind

    Returns:
        ParsedDocument

    Raises:
        DocumentParseError: 解析失败（message 含真实原因）
    """
    file_name = file_name or os.path.basename(file_path)
    ext = get_file_extension(file_name)

    if not os.path.exists(file_path):
        raise DocumentParseError(f"文件不存在: {file_path}")
    if os.path.getsize(file_path) == 0:
        raise DocumentParseError("文件内容为空（0 字节）")

    warnings: List[str] = []

    try:
        # 1) 纯文本类：本地解析，绝不调用外部 API
        if ext in PLAIN_TEXT_EXTENSIONS:
            return ParsedDocument(_read_text_file(file_path), parser="local-text", warnings=warnings)
        if ext in JSON_EXTENSIONS:
            return ParsedDocument(_parse_json(file_path), parser="local-json", warnings=warnings)
        if ext in CSV_EXTENSIONS:
            return ParsedDocument(_parse_csv(file_path), parser="local-csv", warnings=warnings)
        if ext in HTML_EXTENSIONS:
            return ParsedDocument(_parse_html(file_path), parser="local-html", warnings=warnings)

        # 2) 二进制文档类：本地优先，失败再 DocMind
        local_parsers = {
            **{e: ("local-pdf", _parse_pdf_local) for e in PDF_EXTENSIONS},
            **{e: ("local-docx", _parse_docx_local) for e in WORD_EXTENSIONS},
            **{e: ("local-xlsx", _parse_xlsx_local) for e in EXCEL_EXTENSIONS},
            **{e: ("local-pptx", _parse_pptx_local) for e in PPT_EXTENSIONS},
        }

        if ext in local_parsers:
            parser_name, parser_func = local_parsers[ext]
            # .doc / .xls / .ppt 是老二进制格式，本地库不支持，直接走 DocMind
            legacy_binary = ext in {".doc", ".xls", ".ppt"}
            if not legacy_binary:
                try:
                    text = parser_func(file_path)
                    if text and text.strip():
                        return ParsedDocument(text, parser=parser_name, warnings=warnings)
                    warnings.append(f"{parser_name} 未提取到文本（可能是扫描件），尝试 DocMind OCR")
                except DocumentParseError as e:
                    warnings.append(f"{parser_name} 失败: {e}")

            if allow_docmind:
                text = _parse_with_docmind(file_path, file_name)
                if text and text.strip():
                    return ParsedDocument(text, parser="docmind", warnings=warnings)
                raise DocumentParseError(
                    "DocMind 返回内容为空" + (f"；本地解析情况: {'; '.join(warnings)}" if warnings else "")
                )

            raise DocumentParseError(
                "本地解析未获得文本且 DocMind 未启用"
                + (f"；详情: {'; '.join(warnings)}" if warnings else "")
            )

        # 3) 图片：仅 DocMind OCR
        if ext in IMAGE_EXTENSIONS:
            if not allow_docmind:
                raise DocumentParseError("图片文件需要 DocMind OCR，但当前未启用")
            text = _parse_with_docmind(file_path, file_name)
            if not text or not text.strip():
                raise DocumentParseError("图片 OCR 未识别到文本")
            return ParsedDocument(text, parser="docmind-ocr", warnings=warnings)

        # 4) 未知类型：尝试按文本读取
        text = _read_text_file(file_path)
        if text and text.strip():
            warnings.append(f"未知扩展名 {ext or '(空)'}，已按纯文本解析")
            return ParsedDocument(text, parser="local-text-fallback", warnings=warnings)

        raise DocumentParseError(f"不支持的文件类型: {ext or '(无扩展名)'}")

    except DocumentParseError:
        raise
    except Exception as e:
        raise DocumentParseError(f"解析异常: {type(e).__name__}: {e}") from e


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
    """
    文本切分（按句子边界尽量对齐）

    Args:
        text: 原始文本
        chunk_size: 每块目标长度
        overlap: 相邻块重叠长度

    Returns:
        文本块列表
    """
    if not text:
        return []

    # 防御：非法参数会导致死循环
    chunk_size = max(int(chunk_size), 50)
    overlap = min(max(int(overlap), 0), chunk_size // 2)

    chunks: List[str] = []
    start = 0
    text_length = len(text)

    while start < text_length:
        end = min(start + chunk_size, text_length)
        chunk = text[start:end]

        if end < text_length:
            for sep in ["\n\n", "。", "！", "？", ".", "!", "?", "\n"]:
                last_sep = chunk.rfind(sep)
                if last_sep > chunk_size // 2:
                    chunk = chunk[: last_sep + len(sep)]
                    end = start + last_sep + len(sep)
                    break

        if chunk.strip():
            chunks.append(chunk.strip())

        next_start = end - overlap
        # 保证严格前进，避免死循环
        start = next_start if next_start > start else end

    return chunks
